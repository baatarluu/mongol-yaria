#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
mongol-yaria дээр үүсгэсэн Claude Code Skill-үүдийг ашиглах
алхамчилсан гарын авлага (жишээтэй) — PPT үүсгэгч.

Ажиллуулах:
    pip install python-pptx
    python3 docs/build_skills_ppt.py
Гаралт:
    docs/mongol-yaria-skills-guide.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Өнгөний загвар ----
BG      = RGBColor(0x1A, 0x16, 0x12)
CARD    = RGBColor(0x26, 0x21, 0x1B)
DARK    = RGBColor(0x14, 0x11, 0x0D)
ACCENT  = RGBColor(0xD9, 0x77, 0x57)   # улбар шар
ACCENT2 = RGBColor(0xE8, 0xB4, 0x65)   # алтлаг
GREEN   = RGBColor(0x86, 0xC2, 0x32)   # ногоон
BLUE    = RGBColor(0x6A, 0xB0, 0xD9)   # цэнхэр
TEXT    = RGBColor(0xF2, 0xEC, 0xE3)
MUTED   = RGBColor(0xB5, 0xAA, 0x99)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, color=None, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.shadow.inherit = False
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.08):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after)
        p.space_before = Pt(0); p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = "Calibri"
    return tb


def chip(slide, l, t, label, fill, fg=BG, w=Inches(2.6)):
    h = Inches(0.42)
    sp = box(slide, l, t, w, h, color=fill)
    tf = sp.text_frame; tf.word_wrap = False
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(12.5); r.font.bold = True; r.font.color.rgb = fg
    return sp


# ============================================================
# 1) ТАЙТЛ
# ============================================================
s = prs.slides.add_slide(BLANK); bg(s)
box(s, Inches(0), Inches(0), Inches(0.28), SH, color=GREEN)
text(s, Inches(0.9), Inches(1.9), Inches(11.6), Inches(1.4),
     [[("mongol-yaria · Claude Code Skill-үүд", 44, TEXT, True)]])
text(s, Inches(0.95), Inches(3.2), Inches(11.6), Inches(1.0),
     [[("Алхамчилсан ашиглах заавар — жишээтэй", 24, ACCENT, False)]])
text(s, Inches(0.95), Inches(5.9), Inches(11.6), Inches(1.0),
     [[("4 skill: /claude-code-tour · /deploy-helper · /pwa-check · /release-notes", 15, MUTED, False)],
      [("Байршил: .claude/skills/  ·  Эх сурвалж: youtu.be/vfWTyEreOEc", 13, MUTED, False)]])


# ============================================================
# 2) ЮУ ҮҮСГЭСЭН БЭ — OVERVIEW
# ============================================================
s = prs.slides.add_slide(BLANK); bg(s)
text(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.8),
     [[("Юу үүсгэсэн бэ?", 36, TEXT, True)]])
text(s, Inches(0.72), Inches(1.25), Inches(12), Inches(0.6),
     [[("4 ширхэг бэлэн skill — Claude Code дотроос шууд дуудаж ашиглана", 16, MUTED, False)]])

cards = [
    ("/claude-code-tour", GREEN, "Claude Code-ийн 12 боломжийг монголоор зааварлах, дасгал хийлгэх"),
    ("/deploy-helper", ACCENT, "Netlify / Vercel дээр GEMINI_API_KEY-тэй нь deploy хийх"),
    ("/pwa-check", BLUE, "manifest / service worker / icon шалгаж PWA суулгах боломжийг батлах"),
    ("/release-notes", ACCENT2, "git log-оос монголоор changelog / release notes үүсгэх"),
]
x0, y0 = Inches(0.7), Inches(2.1)
cw, ch = Inches(5.95), Inches(2.25)
gx, gy = Inches(0.15), Inches(0.2)
for i, (name, col, desc) in enumerate(cards):
    cx = x0 + (cw + gx) * (i % 2)
    cy = y0 + (ch + gy) * (i // 2)
    box(s, cx, cy, cw, ch, color=CARD)
    box(s, cx, cy, cw, Inches(0.12), color=col)
    text(s, cx + Inches(0.3), cy + Inches(0.35), cw - Inches(0.6), Inches(0.7),
         [[(name, 24, col, True)]])
    text(s, cx + Inches(0.3), cy + Inches(1.2), cw - Inches(0.6), Inches(0.9),
         [[(desc, 15, TEXT, False)]], line_spacing=1.15)


# ============================================================
# 3) SKILL-ИЙГ ЯАЖ АЖИЛЛУУЛАХ ВЭ
# ============================================================
s = prs.slides.add_slide(BLANK); bg(s)
text(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.8),
     [[("Skill-ийг яаж ажиллуулах вэ?", 36, TEXT, True)]])

ways = [
    ("1 · Шууд дуудах", ACCENT,
     "Claude Code-д `/нэр` гэж бичнэ. Жишээ: `/deploy-helper`. "
     "Зарим нь нэмэлт текст авна: `/release-notes сүүлийн tag-аас хойш`."),
    ("2 · Автоматаар", GREEN,
     "Зүгээр л монголоор асуувал тохирох skill өөрөө идэвхжинэ. "
     "Ж: “Энэ аппыг Netlify дээр тавья” → /deploy-helper автоматаар ачаална."),
    ("3 · Хаана байгаа", BLUE,
     ".claude/skills/<нэр>/SKILL.md дотор. Frontmatter-ийн `description` нь "
     "“хэзээ ажиллахыг” тодорхойлдог тул Claude тохирох үед нь сонгоно."),
]
y = Inches(1.55)
for title, col, desc in ways:
    box(s, Inches(0.7), y, Inches(11.9), Inches(1.55), color=CARD)
    box(s, Inches(0.7), y, Inches(0.1), Inches(1.55), color=col)
    text(s, Inches(1.05), y + Inches(0.2), Inches(11.2), Inches(0.5),
         [[(title, 20, col, True)]])
    text(s, Inches(1.05), y + Inches(0.78), Inches(11.2), Inches(0.7),
         [[(desc, 15, TEXT, False)]], line_spacing=1.15)
    y = y + Inches(1.78)

text(s, Inches(0.72), Inches(7.0), Inches(12), Inches(0.4),
     [[("💡 Бүх skill-ийг харах: Claude Code дотор / товчоор жагсаалт гарна", 13, MUTED, False)]])


# ============================================================
# 4) SKILL ДЭЛГЭРЭНГҮЙ СЛАЙД (заавар + жишээ)
# ============================================================
def skill_slide(name, col, what, steps, example, trig):
    s = prs.slides.add_slide(BLANK); bg(s)
    box(s, Inches(0), Inches(0), Inches(0.22), SH, color=col)
    text(s, Inches(0.55), Inches(0.45), Inches(8.5), Inches(0.9),
         [[(name, 32, col, True)]])
    chip(s, Inches(10.0), Inches(0.58), "Skill", col, w=Inches(2.6))
    text(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.7),
         [[("Юу хийдэг:  ", 16, ACCENT2, True), (what, 16, MUTED, False)]],
         line_spacing=1.12)

    # Зүүн: алхамчилсан заавар
    box(s, Inches(0.6), Inches(2.35), Inches(6.55), Inches(4.7), color=CARD)
    text(s, Inches(0.9), Inches(2.55), Inches(6.0), Inches(0.5),
         [[("📋 Алхамчилсан заавар", 18, col, True)]])
    runs = [[(f"{i}. ", 14.5, ACCENT2, True), (st, 14.5, TEXT, False)]
            for i, st in enumerate(steps, 1)]
    text(s, Inches(0.9), Inches(3.15), Inches(6.0), Inches(3.8),
         runs, space_after=8, line_spacing=1.12)

    # Баруун: жишээ
    box(s, Inches(7.35), Inches(2.35), Inches(5.35), Inches(4.7), color=DARK)
    text(s, Inches(7.6), Inches(2.55), Inches(4.9), Inches(0.5),
         [[("💬 Жишээ", 18, GREEN, True)]])
    ex_runs = []
    for who, msg in example:
        if who == "u":
            ex_runs.append([("🗣 ", 14, BLUE, True), (msg, 14, TEXT, True)])
        elif who == "c":
            ex_runs.append([("🤖 ", 14, ACCENT, True), (msg, 13.5, MUTED, False)])
        else:
            ex_runs.append([(msg, 13, GREEN, False)])
    text(s, Inches(7.6), Inches(3.15), Inches(4.9), Inches(3.8),
         ex_runs, space_after=9, line_spacing=1.12)

    text(s, Inches(0.6), Inches(7.12), Inches(12), Inches(0.35),
         [[("🔑 Дуудах түлхүүр үгс:  ", 12.5, ACCENT2, True), (trig, 12.5, MUTED, False)]])
    return s


# --- /claude-code-tour ---
skill_slide(
    "/claude-code-tour", GREEN,
    "Claude Code-ийн 12 боломжийг монголоор тайлбарлаж, энэ repo дээр дасгал хийлгэнэ.",
    [
        "Claude Code дотор `/claude-code-tour` гэж бич.",
        "Түвшингээ сонго: шинэхэн / дунд / ахисан.",
        "Бүх 12 боломжийг дараалан, эсвэл сонгосон 1-2-ыг асуу.",
        "Боломж бүрд: тайлбар + команд + энэ repo дээрх жишээ гарна.",
        "Санал болгосон дасгалыг хийж бататга.",
    ],
    [
        ("u", "/claude-code-tour"),
        ("c", "Түвшингээ сонгоно уу: шинэхэн / дунд / ахисан?"),
        ("u", "Шинэхэн. Subagents-ийг заа."),
        ("c", "Subagents = тусдаа цэвэр контексттэй туслах. `/agents`-ээр нээнэ…"),
        ("o", "Дасгал: app.js-ийг судлах агент үүсгэ →"),
    ],
    "“claude code боломжууд”, “12 features”, “сургалт”, “tour”")

# --- /deploy-helper ---
skill_slide(
    "/deploy-helper", ACCENT,
    "Аппыг Netlify / Vercel дээр HTTPS-ээр байршуулж, AI-г GEMINI_API_KEY-ээр идэвхжүүлнэ.",
    [
        "`/deploy-helper` гэж бичих эсвэл “deploy хийе” гэж асуу.",
        "Платформоо сонго: Netlify эсвэл Vercel.",
        "GitHub-аар нэвтэрч repo-г Import → Deploy хий.",
        "aistudio.google.com/apikey-аас үнэгүй key ав.",
        "Env-д GEMINI_API_KEY нэмж ДАХИН deploy хий.",
        "HTTPS холбоосыг нээж микрофон + AI-г шалга.",
    ],
    [
        ("u", "Netlify дээр тавиад AI-г асаая"),
        ("c", "1) netlify.com→GitHub 2) Import repo 3) Deploy"),
        ("c", "Дараа нь Site settings→Environment variables→"),
        ("o", "GEMINI_API_KEY = <таны key>"),
        ("c", "…нэмээд дахин deploy. AI товч идэвхжинэ ✅"),
    ],
    "“deploy”, “байршуулах”, “Netlify”, “Vercel”, “AI ажиллахгүй байна”")

# --- /pwa-check ---
skill_slide(
    "/pwa-check", BLUE,
    "manifest, service worker, icon, HTTPS-г шалгаж “апп болгож суулгах” боломжийг батална.",
    [
        "`/pwa-check` гэж бич эсвэл “суулгах боломжгүй байна” гэж асуу.",
        "manifest-ийн талбар, icon (192/512 + maskable)-г шалгана.",
        "sw.js бүртгэгдсэн ба APP_SHELL бүрэн эсэхийг хардаг.",
        "HTTPS / DevTools→Application дээрх алдааг илрүүлнэ.",
        "Олдсон асуудал бүрд тодорхой засвар санал болгоно.",
    ],
    [
        ("u", "/pwa-check"),
        ("c", "manifest ✅  icon-192/512 ✅  maskable ✅"),
        ("c", "sw.js бүртгэгдсэн ✅  APP_SHELL бүрэн ✅"),
        ("c", "⚠ Анхаар: шинэ хувилбарт CACHE нэрийг"),
        ("o", "\"mongol-yaria-v1\" → \"v2\" болго"),
    ],
    "“PWA”, “install”, “суулгах боломжгүй”, “офлайн ажиллахгүй”, “manifest”")

# --- /release-notes ---
skill_slide(
    "/release-notes", ACCENT2,
    "git түүхээс хэрэглэгчид зориулсан, монгол хэлээр changelog автоматаар бичнэ.",
    [
        "`/release-notes` гэж бич (хүсвэл хугацаа/tag нэм).",
        "Сүүлийн tag-аас хойших commit-уудыг уншина.",
        "Техникийн биш, ойлгомжтой хэлээр хөрвүүлнэ.",
        "Шинэ боломж / Засвар / Сайжруулалт-аар ангилна.",
        "CHANGELOG.md-д нэмэхийг санал болгоно.",
    ],
    [
        ("u", "/release-notes сүүлийн хувилбараас хойш"),
        ("c", "## v1.1 — 2026-06-08"),
        ("o", "🎤 Шинэ боломж: Япон орчуулга"),
        ("o", "🐛 Засвар: микрофон зогсох алдаа"),
        ("c", "CHANGELOG.md-д нэмэх үү?"),
    ],
    "“release notes”, “changelog”, “хувилбарын тэмдэглэл”, “юу өөрчлөгдсөн”")


# ============================================================
# 8) ХУРААНГУЙ ХҮСНЭГТ
# ============================================================
s = prs.slides.add_slide(BLANK); bg(s)
box(s, Inches(0), Inches(0), Inches(0.28), SH, color=ACCENT)
text(s, Inches(0.9), Inches(0.5), Inches(11.5), Inches(0.9),
     [[("Хураангуй — түргэн лавлах", 34, TEXT, True)]])

rows = [
    ("/claude-code-tour", "Сурах / заах", "/claude-code-tour", GREEN),
    ("/deploy-helper", "Production-д гаргах", "“Netlify дээр deploy хийе”", ACCENT),
    ("/pwa-check", "Суулгах боломж батлах", "/pwa-check", BLUE),
    ("/release-notes", "Хувилбар гаргах", "/release-notes", ACCENT2),
]
# толгой
hy = Inches(1.55)
for cx, w, label in [(Inches(0.9), Inches(3.6), "Skill"),
                     (Inches(4.6), Inches(3.6), "Хэзээ ашиглах"),
                     (Inches(8.3), Inches(4.1), "Жишээ дуудлага")]:
    text(s, cx, hy, w, Inches(0.5), [[(label, 15, MUTED, True)]])
y = Inches(2.1)
for name, when, ex, col in rows:
    box(s, Inches(0.9), y, Inches(11.5), Inches(0.92), color=CARD)
    box(s, Inches(0.9), y, Inches(0.1), Inches(0.92), color=col)
    text(s, Inches(1.2), y + Inches(0.22), Inches(3.4), Inches(0.5),
         [[(name, 16, col, True)]])
    text(s, Inches(4.6), y + Inches(0.24), Inches(3.6), Inches(0.5),
         [[(when, 15, TEXT, False)]])
    text(s, Inches(8.3), y + Inches(0.24), Inches(4.1), Inches(0.5),
         [[(ex, 14, GREEN, False)]])
    y = y + Inches(1.07)

text(s, Inches(0.9), Inches(6.95), Inches(11.5), Inches(0.4),
     [[("Дэлгэрэнгүй: docs/README.md  ·  12 боломжийн танилцуулга: docs/claude-code-12-features.pptx", 12.5, MUTED, False)]])


prs.save("docs/mongol-yaria-skills-guide.pptx")
print("OK -> docs/mongol-yaria-skills-guide.pptx  (slides:", len(prs.slides._sldIdLst), ")")
