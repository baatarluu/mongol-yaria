#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Claude Code-ийн 12 шилдэг боломж — алхамчилсан PPT үүсгэгч.

Эх сурвалж: Nate Herk, "I Tested Every Claude Code Feature, These 12 Are the Best"
            https://youtu.be/vfWTyEreOEc

Ажиллуулах:
    pip install python-pptx
    python3 docs/build_ppt.py
Гаралт:
    docs/claude-code-12-features.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Өнгөний загвар (Claude / тагтаа улбар шар онцлох) ----
BG       = RGBColor(0x1A, 0x16, 0x12)   # бараан фон
CARD     = RGBColor(0x26, 0x21, 0x1B)   # карт
ACCENT   = RGBColor(0xD9, 0x77, 0x57)   # Claude улбар шар
ACCENT2  = RGBColor(0xE8, 0xB4, 0x65)   # алтлаг
TEXT     = RGBColor(0xF2, 0xEC, 0xE3)   # цайвар текст
MUTED    = RGBColor(0xB5, 0xAA, 0x99)   # бүдэг текст
GREEN    = RGBColor(0x86, 0xC2, 0x32)   # ногоон онцлох

SW, SH = Inches(13.333), Inches(7.5)   # 16:9

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, color=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.shadow.inherit = False
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    return sp


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.05):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold)."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def chip(slide, l, t, label, fill, fg=BG):
    w = Inches(0.95)
    h = Inches(0.42)
    sp = box(slide, l, t, w, h, color=fill)
    tf = sp.text_frame
    tf.word_wrap = False
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = fg
    return sp


# ============================================================
# 1) ТАЙТЛ СЛАЙД
# ============================================================
s = prs.slides.add_slide(BLANK); bg(s)
box(s, Inches(0), Inches(0), Inches(0.28), SH, color=ACCENT)
text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2), [
    [("Claude Code", 54, TEXT, True), ("  —  12 шилдэг боломж", 54, ACCENT, True)],
])
text(s, Inches(0.95), Inches(3.25), Inches(11.5), Inches(1.0), [
    [("Тэдгээрийг ашиглах алхамчилсан гарын авлага", 24, MUTED, False)],
])
text(s, Inches(0.95), Inches(5.7), Inches(11.5), Inches(1.2), [
    [("Эх сурвалж: Nate Herk — “I Tested Every Claude Code Feature, These 12 Are the Best”", 14, MUTED, False)],
    [("youtu.be/vfWTyEreOEc", 14, ACCENT2, False)],
])


# ============================================================
# 2) АГУУЛГА / ТИР ЖАГСААЛТ
# ============================================================
s = prs.slides.add_slide(BLANK); bg(s)
text(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.8),
     [[("Агуулга — Tier эрэмбэ", 36, TEXT, True)]])

tiers = [
    ("S-Tier · Хамгийн өндөр нөлөөтэй", ACCENT, [
        "1. CLAUDE.md — Төслийн санах ой",
        "2. Subagents — Зэрэгцээ дэд агентууд",
        "3. Custom Commands — Тохируулсан ажлын урсгал",
    ]),
    ("A-Tier · Өдөр тутмын зайлшгүй", ACCENT2, [
        "4. Plan Mode — Төлөвлөх горим",
        "5. Context Lifecycle — /clear, /compact",
        "6. Rollback Checkpoints — /rewind",
        "7. MCP — Гадаад хэрэгсэл холбох",
        "8. Hooks — Автомат хяналт",
    ]),
    ("A-Tier · Өргөтгөл ба бүтээмж", GREEN, [
        "9. Skills — Дахин ашиглах чадвар",
        "10. Web Search & Fetch",
        "11. File / Image Upload",
        "12. Task Scheduling & IDE",
    ]),
]
x = Inches(0.7)
colw = Inches(3.95)
gap = Inches(0.13)
for title, col, items in tiers:
    box(s, x, Inches(1.45), colw, Inches(5.4), color=CARD)
    box(s, x, Inches(1.45), colw, Inches(0.12), color=col)
    text(s, x + Inches(0.25), Inches(1.7), colw - Inches(0.5), Inches(0.9),
         [[(title, 16, col, True)]])
    runs = [[("•  " + it, 15, TEXT, False)] for it in items]
    text(s, x + Inches(0.25), Inches(2.75), colw - Inches(0.5), Inches(4.0),
         runs, space_after=10, line_spacing=1.1)
    x = x + colw + gap


# ============================================================
# 3) ФИЧЕР СЛАЙДУУД (алхам алхмаар)
# ============================================================
def feature_slide(num, tier, tier_col, title, what, steps, tip, cmds=None):
    s = prs.slides.add_slide(BLANK); bg(s)
    box(s, Inches(0), Inches(0), Inches(0.22), SH, color=tier_col)
    # дугаар
    text(s, Inches(0.55), Inches(0.4), Inches(1.6), Inches(1.2),
         [[(f"{num:02d}", 44, tier_col, True)]])
    # tier chip
    chip(s, Inches(11.4), Inches(0.55), tier, tier_col)
    # гарчиг
    text(s, Inches(1.7), Inches(0.5), Inches(9.5), Inches(1.0),
         [[(title, 32, TEXT, True)]])
    # "Юу вэ?"
    text(s, Inches(1.7), Inches(1.45), Inches(10.9), Inches(0.9),
         [[("Юу вэ?  ", 16, ACCENT2, True), (what, 16, MUTED, False)]],
         line_spacing=1.1)

    # Алхмууд карт (зүүн)
    box(s, Inches(0.7), Inches(2.5), Inches(7.3), Inches(4.5), color=CARD)
    text(s, Inches(1.0), Inches(2.7), Inches(6.8), Inches(0.5),
         [[("Алхамчилсан хэрэглээ", 18, ACCENT, True)]])
    runs = []
    for i, st in enumerate(steps, 1):
        runs.append([(f"{i}. ", 15, ACCENT2, True), (st, 15, TEXT, False)])
    text(s, Inches(1.0), Inches(3.3), Inches(6.8), Inches(3.6),
         runs, space_after=9, line_spacing=1.12)

    # Командууд + Tip карт (баруун)
    box(s, Inches(8.25), Inches(2.5), Inches(4.4), Inches(2.6), color=RGBColor(0x14,0x11,0x0D))
    text(s, Inches(8.5), Inches(2.7), Inches(3.9), Inches(0.5),
         [[("⌨  Команд / Байршил", 15, ACCENT2, True)]])
    cmd_runs = [[(c, 14, GREEN, False)] for c in (cmds or ["—"])]
    text(s, Inches(8.5), Inches(3.25), Inches(3.9), Inches(1.7),
         cmd_runs, space_after=6, line_spacing=1.1)

    box(s, Inches(8.25), Inches(5.25), Inches(4.4), Inches(1.75), color=CARD)
    text(s, Inches(8.5), Inches(5.42), Inches(3.9), Inches(1.5),
         [[("💡 Зөвлөмж", 15, ACCENT, True)],
          [(tip, 13.5, MUTED, False)]],
         space_after=4, line_spacing=1.1)
    return s


S = ("S-Tier", ACCENT)
A = ("A-Tier", ACCENT2)
A3 = ("A-Tier", GREEN)

feature_slide(1, *S,
    "CLAUDE.md — Төслийн санах ой",
    "Claude нэг бүр session-д уншдаг “үндсэн хууль”. Төслийн дүрэм, командуудаа дахин тайлбарлахгүйгээр санадаг.",
    [
        "Терминал дээр `claude` нээгээд `/init` гэж бичнэ → стартер CLAUDE.md үүснэ.",
        "Доторх дүрмийг засна: ашигладаг командууд, код стиль, хориглосон зүйлс.",
        "`/memory` командаар хүссэн үедээ нэмж сайжруулна.",
        "Хувийн (глобал) санах ойг `~/.claude/CLAUDE.md`-д бичнэ.",
        "Чухал зүйлийг # тэмдгээр сануулбал автоматаар хадгалагдана.",
    ],
    "Богино, тодорхой байлга. 20-40 мөр хангалттай — урт CLAUDE.md context-ийг дэмий иднэ.",
    ["/init", "/memory", "CLAUDE.md", "~/.claude/CLAUDE.md"])

feature_slide(2, *S,
    "Subagents — Зэрэгцээ дэд агентууд",
    "Тусдаа цэвэр контексттэй дэд агент үүсгэж, судалгаа/шалгалтыг даалгаад зөвхөн дүгнэлтийг нь авна.",
    [
        "`/agents` командаар агент менежерийг нээнэ.",
        "Шинэ агент үүсгэж нэр, зорилго, ашиглах хэрэгслийг тодорхойлно.",
        "Эсвэл `.claude/agents/<нэр>.md` файл үүсгэж frontmatter бичнэ.",
        "Том ажлыг хэд хэдэн агентад зэрэг хуваарилбал параллель ажиллана.",
        "Үндсэн контекст цэвэр хэвээр үлдэж, зөвхөн үр дүн буцна.",
    ],
    "Судалгаа, код ревью, тест зэрэг “их уншдаг” ажлыг агентад өг — гол контекстоо хэмнэнэ.",
    ["/agents", "/tasks", ".claude/agents/"])

feature_slide(3, *S,
    "Custom Commands — Ажлын урсгал",
    "Байнга давтагддаг promt-уудыг `/команд` болгож хадгална. Багаар нэг стандарт урсгал хуваалцана.",
    [
        "`.claude/commands/` хавтас үүсгэнэ.",
        "Доторх `deploy.md` гэх мэт файл = `/deploy` команд.",
        "Файлд хийх ажлын зааврыг бичнэ ($ARGUMENTS-аар параметр авна).",
        "Төсөлд хадгалбал багийн гишүүд бүгд адил командтай болно.",
        "Илүү ухаалаг хувилбар нь Skills (9-р боломжийг үз).",
    ],
    "Repo-д commit хийвэл бүх хүн нэг л `/команд`-аар адилхан үр дүн авна.",
    [".claude/commands/", "deploy.md → /deploy", "$ARGUMENTS"])

feature_slide(4, *A,
    "Plan Mode — Төлөвлөх горим",
    "Том өөрчлөлт хийхээс өмнө Claude кодыг засахгүйгээр зөвхөн төлөвлөгөө гаргадаг горим.",
    [
        "Shift+Tab дарж горим солих (эсвэл plan mode-руу шилжих).",
        "Хүсэлтээ бичнэ — Claude кодыг УНШИНА, гэхдээ ЗАСАХГҮЙ.",
        "Гаргасан төлөвлөгөөг уншиж шалгана.",
        "Зөвшөөрвөл горимоос гарч хэрэгжүүлэлт эхэлнэ.",
        "Эрсдэлтэй/том рефактор бүрт эхлээд Plan Mode ашигла.",
    ],
    "“Эхлээд төлөвлө, дараа нь хэрэгжүүл” — алдаатай засваруудаас сэргийлнэ.",
    ["Shift + Tab", "(plan mode)"])

feature_slide(5, *A,
    "Context Lifecycle — /clear, /compact",
    "Урт session-ы контекстийг цэвэрлэх буюу хураангуйлж, Claude-ыг анхааралтай, хурдан байлгана.",
    [
        "Шинэ ажил эхлэхэд `/clear` — контекстийг бүрэн цэвэрлэнэ.",
        "Урт яриаг үргэлжлүүлэхдээ `/compact` — хураангуйлж шахна.",
        "Контекст дүүрэх үед чанар буурдаг тул тогтмол цэвэрлэ.",
        "Сэдэв солих бүрт `/clear` дарж зуршил болго.",
    ],
    "Нэг таб дотор олон сэдэв холих нь алдаа дагуулдаг — сэдэв бүрд цэвэр контекст.",
    ["/clear", "/compact"])

feature_slide(6, *A,
    "Rollback Checkpoints — /rewind",
    "Claude prompt бүр дээр автомат checkpoint үүсгэдэг тул алдаатай өөрчлөлтийг буцаах боломжтой.",
    [
        "Хүсэлт бүрд checkpoint автоматаар үүснэ.",
        "Буцаахын тулд Esc товчийг хоёр дарна (Esc, Esc).",
        "Эсвэл `/rewind` командаар цэгийг сонгож буцна.",
        "Кодын болон ярианы төлөвийг хоёуланг нь сэргээж болно.",
    ],
    "“Болохгүй боллоо” гэхээс бүү ай — /rewind-ээр аюулгүй туршиж үз.",
    ["/rewind", "Esc Esc"])

feature_slide(7, *A,
    "MCP — Гадаад хэрэгсэл холбох",
    "Model Context Protocol: GitHub, өгөгдлийн сан, браузер, дотоод API зэргийг Claude-д холбоно.",
    [
        "`/mcp` командаар холбогдсон сервэрүүдийг харна.",
        "MCP сервэрийг тохиргоонд нэмнэ (нэр, команд, орчны хувьсагч).",
        "Жишээ: GitHub MCP-ээр PR, issue, CI-г шууд удирдана.",
        "3000+ интеграц байдаг — Sentry, Slack, Postgres г.м.",
        "Claude сервэрийн өгөгдлийг ашиглан логикийг боддог.",
    ],
    "Эхлээд хэрэгтэй 1-2 сервэрээ л холбо — хэт олон хэрэгсэл контекстийг дүүргэнэ.",
    ["/mcp", ".mcp.json", "claude mcp add"])

feature_slide(8, *A,
    "Hooks — Автомат хяналт",
    "Тодорхой үйлдлийн өмнө/дараа автоматаар ажиллах скрипт. Чанар, аюулгүй байдлын “хамгаалалт”.",
    [
        "`.claude/settings.json`-д hooks тохируулна.",
        "PreToolUse — хэрэгсэл ажиллахаас ӨМНӨ (зөвшөөрөл шалгах).",
        "PostToolUse — ДАРАА нь (жишээ: TypeScript шалгах, formatter).",
        "UserPromptSubmit, Stop, Notification гэх мэт олон үйлдэл бий.",
        "Жишээ: засвар бүрийн дараа `npm run lint` автоматаар ажиллуулах.",
    ],
    "Hook бол Claude биш, харин harness өөрөө ажиллуулдаг — найдвартай автоматжуулалт.",
    ["PreToolUse", "PostToolUse", ".claude/settings.json"])

feature_slide(9, *A3,
    "Skills — Дахин ашиглах чадвар",
    "SKILL.md файл болгон савласан тусгай чадвар. Claude даалгаварт тохирохыг нь автоматаар сонгож ачаална.",
    [
        "`.claude/skills/<нэр>/SKILL.md` үүсгэнэ.",
        "Frontmatter-т name + description заавал бичнэ.",
        "description-д хэзээ ашиглахыг тодорхой бич — Claude үүгээр сонгоно.",
        "Дэлгэрэнгүй заавар, скрипт, файлуудыг хавтаст нэмж болно.",
        "`/<нэр>` гэж бичээд эсвэл автоматаар дуудагдана.",
    ],
    "Progressive disclosure: зөвхөн хэрэгтэй skill л ачаалагдана — контекстод хэмнэлттэй.",
    [".claude/skills/<нэр>/SKILL.md", "name / description"])

feature_slide(10, *A3,
    "Web Search & Fetch",
    "Claude интернэтээс шинэ мэдээлэл хайж, тодорхой URL-ийн агуулгыг татаж ажилдаа ашиглана.",
    [
        "Зүгээр л “…-ийг вэбээс хайж үз” гэж асууна.",
        "Сүүлийн үеийн мэдээлэл, баримт бичиг, API өөрчлөлт татна.",
        "Тодорхой холбоосын агуулгыг шууд уншуулж болно.",
        "Хайлтын дараа эх сурвалжаа жагсааж өгдөг.",
    ],
    "Сургалтын мэдлэгийн огнооноос хойшхи зүйлийг асуухдаа вэб хайлт ашиглуул.",
    ["WebSearch", "WebFetch"])

feature_slide(11, *A3,
    "File / Image Upload",
    "Зураг, дизайн, screenshot, лог зэргийг чатад өгч, Claude-аар шинжлүүлж, кодлуулна.",
    [
        "Файл/зургийг чат руу чирж оруулна (drag & drop).",
        "Screenshot өгч UI-г кодлуулах (“энэ дизайныг хий”).",
        "Алдааны лог, error screenshot өгч дебаг хийлгэх.",
        "PDF, баримт уншуулж агуулгыг боловсруулах.",
    ],
    "Урт алдааны мессеж бичихээс илүү screenshot өгөх нь хурдан бөгөөд тодорхой.",
    ["Drag & drop", "Paste image"])

feature_slide(12, *A3,
    "Task Scheduling & IDE",
    "Удаан ажлыг background-д төлөвлөж ажиллуулах, мөн VS Code / JetBrains дотроос шууд ашиглах.",
    [
        "Background task-аар урт ажлыг зэрэг ажиллуулна.",
        "`/tasks`-аар ажиллаж буй ажлуудыг харна.",
        "VS Code / JetBrains extension суулгаж IDE дотроо ашиглана.",
        "Cloud / web (claude.ai/code), утсан дээрх app-аар үргэлжлүүлнэ.",
    ],
    "Хүнд build/тестийг background-д тавиад, гол ажлаа үргэлжлүүл.",
    ["/tasks", "VS Code / JetBrains", "claude.ai/code"])


# ============================================================
# 4) ЭНЭ REPO-Д БЭЛЭН СУУЛГАСАН SKILL-УУД
# ============================================================
s = prs.slides.add_slide(BLANK); bg(s)
text(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.8),
     [[("Энэ repo-д бэлэн суулгасан Skill-ууд", 32, TEXT, True)]])
text(s, Inches(0.72), Inches(1.25), Inches(12), Inches(0.6),
     [[("`.claude/skills/` дотор — шууд `/нэр` гэж дуудаж ашиглана", 16, MUTED, False)]])

skills = [
    ("/claude-code-tour", "Энэ 12 боломжийг монголоор зааварлах, дасгал хийлгэх интерактив гарын авлага."),
    ("/deploy-helper", "mongol-yaria-г Netlify эсвэл Vercel дээр GEMINI_API_KEY-тэй нь алхам алхмаар deploy хийнэ."),
    ("/pwa-check", "manifest, service worker, icon-уудыг шалгаж PWA суулгах боломжтой эсэхийг баталгаажуулна."),
    ("/release-notes", "git log-оос монгол хэлээр товч, ойлгомжтой release notes автоматаар үүсгэнэ."),
]
y = Inches(2.05)
for name, desc in skills:
    box(s, Inches(0.7), y, Inches(11.9), Inches(1.05), color=CARD)
    box(s, Inches(0.7), y, Inches(0.1), Inches(1.05), color=GREEN)
    text(s, Inches(1.0), y + Inches(0.12), Inches(3.3), Inches(0.8),
         [[(name, 18, ACCENT2, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.3), y + Inches(0.12), Inches(8.0), Inches(0.85),
         [[(desc, 14.5, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    y = y + Inches(1.2)


# ============================================================
# 5) ХУРААНГУЙ / WORKFLOW
# ============================================================
s = prs.slides.add_slide(BLANK); bg(s)
box(s, Inches(0), Inches(0), Inches(0.28), SH, color=ACCENT)
text(s, Inches(0.9), Inches(0.6), Inches(11.5), Inches(0.9),
     [[("Дүгнэлт — Бүгдийг нэгтгэсэн ажлын урсгал", 32, TEXT, True)]])

flow = [
    ("1 · Тохируулга", "/init → CLAUDE.md, /mcp, /agents, hooks тохируул", ACCENT),
    ("2 · Төлөвлө", "Plan Mode (Shift+Tab) — эхлээд төлөвлөгөө гарга", ACCENT2),
    ("3 · Судал", "Subagents-аар судалгаа/ревью хийлгэ, контекст цэвэр байлга", GREEN),
    ("4 · Хэрэгжүүл", "Skills + Custom commands-аар стандарт урсгалаар ажилла", ACCENT),
    ("5 · Хамгаал", "Hooks-оор lint/test, /rewind-ээр алдаа буцаа", ACCENT2),
    ("6 · Цэвэрлэ", "/clear, /compact — дараагийн ажилд бэлд", GREEN),
]
y = Inches(1.7)
for title, desc, col in flow:
    box(s, Inches(0.9), y, Inches(11.5), Inches(0.78), color=CARD)
    box(s, Inches(0.9), y, Inches(0.1), Inches(0.78), color=col)
    text(s, Inches(1.2), y + Inches(0.08), Inches(3.2), Inches(0.62),
         [[(title, 17, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(4.3), y + Inches(0.08), Inches(8.0), Inches(0.62),
         [[(desc, 15, TEXT, False)]], anchor=MSO_ANCHOR.MIDDLE)
    y = y + Inches(0.9)

text(s, Inches(0.9), Inches(7.05), Inches(11.5), Inches(0.4),
     [[("Эх сурвалж: youtu.be/vfWTyEreOEc · Баримтжуулалт: code.claude.com/docs", 12, MUTED, False)]])


prs.save("docs/claude-code-12-features.pptx")
print("OK -> docs/claude-code-12-features.pptx  (slides:", len(prs.slides._sldIdLst), ")")
